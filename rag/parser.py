"""Document Extraction Module for RAG Pipeline.

Parses PDFs, Word files, Multi-sheet Excel tables, PowerPoint presentations, CSVs, Markdown, JSON, TXT, and Images.
"""

import io
import json
import logging
from typing import Any, Dict, List, Set, Tuple
import pandas as pd

from config import MAX_TABLE_ROWS
from rag.ocr import extract_text_from_image

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS: Set[str] = {
    "pdf", "docx", "pptx", "txt", "md", "json", "csv", "xlsx", "xls", "png", "jpg", "jpeg"
}


def _flatten_json(data: Any, prefix: str = "") -> List[str]:
    """Flattens nested JSON objects/lists into key-value text lines for improved RAG retrieval."""
    lines = []
    if isinstance(data, dict):
        for key, val in data.items():
            new_key = f"{prefix}.{key}" if prefix else str(key)
            lines.extend(_flatten_json(val, new_key))
    elif isinstance(data, list):
        for idx, item in enumerate(data):
            new_key = f"{prefix}[{idx}]"
            lines.extend(_flatten_json(item, new_key))
    else:
        lines.append(f"{prefix}: {data}")
    return lines


def extract_pdf(file_bytes: bytes, filename: str) -> Tuple[List[Dict[str, Any]], str, Dict[str, Any], List[str]]:
    pages_data: List[Dict[str, Any]] = []
    full_text_segments: List[str] = []
    warnings: List[str] = []
    pdfplumber_success = False

    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            if pdf.is_encrypted:
                warnings.append("PDF is encrypted. Extraction may be partial.")
            for page_idx, page in enumerate(pdf.pages):
                txt = page.extract_text() or ""
                pages_data.append({"page_number": page_idx + 1, "text": txt})
                if txt.strip():
                    full_text_segments.append(f"--- Page {page_idx + 1} ---\n{txt}")
                else:
                    warnings.append(f"Page {page_idx + 1} contained no extractable text.")
            pdfplumber_success = True
    except ImportError as imp_err:
        logger.warning(f"pdfplumber not installed for {filename}: {imp_err}. Attempting PyPDF fallback.")
        warnings.append("pdfplumber library missing. Using PyPDF fallback.")
    except Exception as err:
        logger.warning(f"pdfplumber execution failed for {filename}: {err}. Attempting PyPDF fallback.")
        warnings.append(f"Primary PDF parser failed ({type(err).__name__}). Using PyPDF fallback.")

    if not pdfplumber_success:
        pages_data.clear()
        full_text_segments.clear()
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes), strict=False)
            if reader.is_encrypted:
                try:
                    reader.decrypt("")
                except Exception:
                    warnings.append("Could not decrypt protected PDF file.")

            for page_idx, page in enumerate(reader.pages):
                try:
                    txt = page.extract_text() or ""
                    pages_data.append({"page_number": page_idx + 1, "text": txt})
                    if txt.strip():
                        full_text_segments.append(f"--- Page {page_idx + 1} ---\n{txt}")
                except Exception as p_err:
                    warnings.append(f"Failed to read page {page_idx + 1}: {str(p_err)}")
        except Exception as pypdf_err:
            logger.error(f"PyPDF fallback failed for {filename}: {pypdf_err}")
            warnings.append(f"Fatal PDF parsing error: {str(pypdf_err)}")

    metadata = {"total_pages": len(pages_data)}
    return pages_data, "\n\n".join(full_text_segments), metadata, warnings


def extract_docx(file_bytes: bytes, filename: str) -> Tuple[List[Dict[str, Any]], str, Dict[str, Any], List[str]]:
    warnings: List[str] = []
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        joined_text = "\n".join(paragraphs)
        
        pages_data = [{"page_number": 1, "text": joined_text}]
        metadata = {"total_paragraphs": len(paragraphs)}
        return pages_data, joined_text, metadata, warnings
    except Exception as err:
        logger.error(f"DOCX extraction failed for {filename}: {err}")
        warnings.append(f"Failed to parse Word document: {str(err)}")
        return [{"page_number": 1, "text": ""}], "", {"total_paragraphs": 0}, warnings


def extract_pptx(file_bytes: bytes, filename: str) -> Tuple[List[Dict[str, Any]], str, Dict[str, Any], List[str]]:
    warnings: List[str] = []
    pages_data: List[Dict[str, Any]] = []
    full_text_segments: List[str] = []

    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        for idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            slide_content = "\n".join(slide_texts)
            pages_data.append({"page_number": idx + 1, "text": slide_content})
            if slide_content:
                full_text_segments.append(f"--- Slide {idx + 1} ---\n{slide_content}")

    except ImportError:
        warnings.append("python-pptx library not installed. Presentation parsing skipped.")
    except Exception as err:
        logger.error(f"PPTX extraction failed for {filename}: {err}")
        warnings.append(f"Failed to parse PowerPoint presentation: {str(err)}")

    metadata = {"total_slides": len(pages_data)}
    return pages_data, "\n\n".join(full_text_segments), metadata, warnings


def extract_csv_excel(file_bytes: bytes, filename: str, ext: str) -> Tuple[List[Dict[str, Any]], str, Dict[str, Any], List[str]]:
    warnings: List[str] = []
    pages_data: List[Dict[str, Any]] = []
    full_text_segments: List[str] = []
    total_rows = 0

    try:
        if ext == "csv":
            df = None
            for encoding in ["utf-8-sig", "utf-8", "cp1252", "latin-1"]:
                try:
                    # Optimized streaming read: Stops at MAX_TABLE_ROWS directly in memory
                    df = pd.read_csv(
                        io.BytesIO(file_bytes),
                        encoding=encoding,
                        on_bad_lines="skip",
                        nrows=MAX_TABLE_ROWS,
                    )
                    break
                except (UnicodeDecodeError, pd.errors.ParserError, pd.errors.EmptyDataError):
                    continue

            if df is None:
                raise ValueError("CSV decoding failed across all standard encodings and parsers.")

            text_rep = df.to_string(index=False)
            pages_data.append({"page_number": 1, "sheet_name": "CSV", "text": text_rep})
            full_text_segments.append(text_rep)
            total_rows = len(df)
            sheet_names = ["CSV"]

            if total_rows >= MAX_TABLE_ROWS:
                warnings.append(f"CSV read limited to first {MAX_TABLE_ROWS} rows for memory safety.")

        else:
            excel_file = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None, engine=None)
            sheet_names = list(excel_file.keys())
            
            for idx, (sheet, df) in enumerate(excel_file.items()):
                if df.empty:
                    warnings.append(f"Sheet '{sheet}' is empty.")
                    continue
                if len(df) > MAX_TABLE_ROWS:
                    warnings.append(f"Sheet '{sheet}' truncated to first {MAX_TABLE_ROWS} rows.")
                    df = df.head(MAX_TABLE_ROWS)
                
                sheet_text = f"=== Sheet: {sheet} ===\n" + df.to_string(index=False)
                pages_data.append({"page_number": idx + 1, "sheet_name": sheet, "text": sheet_text})
                full_text_segments.append(sheet_text)
                total_rows += len(df)

    except Exception as err:
        logger.error(f"Tabular extraction failed for {filename}: {err}")
        warnings.append(f"Failed to parse tabular file: {str(err)}")
        sheet_names = []

    metadata = {"total_rows": total_rows, "sheets": sheet_names}
    return pages_data, "\n\n".join(full_text_segments), metadata, warnings


def extract_txt_md_json(file_bytes: bytes, filename: str, ext: str) -> Tuple[List[Dict[str, Any]], str, Dict[str, Any], List[str]]:
    warnings: List[str] = []
    raw_text = ""
    
    for encoding in ["utf-8-sig", "utf-8", "cp1252", "latin-1"]:
        try:
            raw_text = file_bytes.decode(encoding)
            break
        except UnicodeDecodeError:
            continue

    if ext == "json":
        try:
            parsed = json.loads(raw_text)
            flattened_lines = _flatten_json(parsed)
            text_content = "\n".join(flattened_lines)
        except Exception:
            warnings.append("JSON structural parse failed; imported as raw text.")
            text_content = raw_text
    else:
        text_content = raw_text

    pages_data = [{"page_number": 1, "text": text_content}]
    metadata = {"character_count": len(text_content)}
    return pages_data, text_content, metadata, warnings


def extract_file(file_bytes: bytes, filename: str) -> Dict[str, Any]:
    if not filename or "." not in filename:
        raise ValueError(f"Invalid filename: '{filename}'. Extension missing.")

    ext = filename.rsplit(".", 1)[-1].lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file format '.{ext}'. Supported extensions: {sorted(list(SUPPORTED_EXTENSIONS))}"
        )

    file_size = len(file_bytes)

    if ext == "pdf":
        pages, text, specific_meta, warnings = extract_pdf(file_bytes, filename)
    elif ext == "docx":
        pages, text, specific_meta, warnings = extract_docx(file_bytes, filename)
    elif ext == "pptx":
        pages, text, specific_meta, warnings = extract_pptx(file_bytes, filename)
    elif ext in ["csv", "xlsx", "xls"]:
        pages, text, specific_meta, warnings = extract_csv_excel(file_bytes, filename, ext)
    elif ext in ["png", "jpg", "jpeg"]:
        img_text, specific_meta, warnings = extract_text_from_image(file_bytes, filename)
        pages = [{"page_number": 1, "text": img_text}]
        text = img_text
    elif ext in ["txt", "md", "json"]:
        pages, text, specific_meta, warnings = extract_txt_md_json(file_bytes, filename, ext)

    metadata = {
        "filename": filename,
        "extension": ext,
        "size_bytes": file_size,
        **specific_meta,
    }

    return {
        "filename": filename,
        "file_type": ext,
        "pages": pages,
        "text": text,
        "metadata": metadata,
        "warnings": warnings,
    }
